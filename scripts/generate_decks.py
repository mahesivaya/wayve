#!/usr/bin/env python3
"""Generate technical_wayve.pptx + business_wayve.pptx from the markdown
deep-dive companion docs.

Both decks use the same visual language (dark hero, blue accent, monospace
for code) so they feel like a set. The slide content is hand-curated —
the markdown docs are too dense for slide-by-slide conversion, so the
decks pick the talk-track points an engineer/board would actually want
to project.

Run from the repo root:
    python3 scripts/generate_decks.py

Output:
    technical_wayve.pptx
    business_wayve.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Visual constants ──────────────────────────────────────────────────

BG_HERO = RGBColor(0x0F, 0x17, 0x2A)      # slate-900
BG_NAV = RGBColor(0x1E, 0x3A, 0x8A)       # blue-900
BG_PANEL = RGBColor(0xFF, 0xFF, 0xFF)
BG_PANEL_DARK = RGBColor(0x1F, 0x29, 0x37)
ACCENT = RGBColor(0x25, 0x63, 0xEB)       # blue-600
ACCENT_LIGHT = RGBColor(0xE0, 0xE7, 0xFF)
INK_PRIMARY = RGBColor(0x11, 0x18, 0x27)
INK_INVERT = RGBColor(0xF9, 0xFA, 0xFB)
INK_MUTED = RGBColor(0x6B, 0x72, 0x80)
INK_MUTED_INVERT = RGBColor(0xCB, 0xD5, 0xE1)
DIVIDER = RGBColor(0xE5, 0xE7, 0xEB)
OK = RGBColor(0x16, 0x65, 0x34)
WARN = RGBColor(0x92, 0x40, 0x0E)
ERR = RGBColor(0x99, 0x1B, 0x1B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────


def set_bg(slide, color):
    """Paint the whole slide background."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # Send to back by removing then re-inserting at index 2 (after layout shapes).
    # python-pptx doesn't expose z-order directly, but inserting first means
    # subsequent shapes render on top of this rectangle, which is what we want.
    return bg


def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=INK_PRIMARY,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=18, color=INK_PRIMARY,
                bullet_color=ACCENT, indent=0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = indent
        # Custom bullet via a colored "•" run + text run.
        bullet = p.add_run()
        bullet.text = "▸  "
        bullet.font.size = Pt(size)
        bullet.font.bold = True
        bullet.font.color.rgb = bullet_color
        body = p.add_run()
        body.text = item
        body.font.size = Pt(size)
        body.font.color.rgb = color
        body.font.name = "Calibri"
    return tb


def add_rect(slide, x, y, w, h, fill, *, line=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line:
        rect.line.color.rgb = line
        rect.line.width = Pt(1)
    else:
        rect.line.fill.background()
    rect.shadow.inherit = False
    return rect


def add_pill(slide, x, y, w, h, text, fill, ink):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = ink
    return pill


def add_table(slide, x, y, w, h, headers, rows, *, header_fill=BG_NAV,
              header_ink=INK_INVERT, row_fill=BG_PANEL, alt_fill=DIVIDER,
              body_ink=INK_PRIMARY, font_size=12):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    tbl = tbl_shape.table
    for c, col_w in enumerate(_col_widths(w, n_cols)):
        tbl.columns[c].width = col_w
    for c, header in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = header
        r.font.size = Pt(font_size + 1)
        r.font.bold = True
        r.font.color.rgb = header_ink
    for ri, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = tbl.cell(ri + 1, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_fill if ri % 2 == 0 else alt_fill
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(value)
            r.font.size = Pt(font_size)
            r.font.color.rgb = body_ink
    return tbl_shape


def _col_widths(total, n):
    base = total // n
    rem = total - base * n
    widths = [base] * n
    widths[0] += rem
    return widths


def add_footer(slide, page_num, total, deck_name):
    add_text(slide, Inches(0.4), Inches(7.05), Inches(8), Inches(0.4),
             f"{deck_name}  ·  Wayve  ·  2026.05",
             size=10, color=INK_MUTED)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.5), Inches(0.4),
             f"{page_num} / {total}",
             size=10, color=INK_MUTED, align=PP_ALIGN.RIGHT)


# ── Slide layout primitives ───────────────────────────────────────────


def make_title_slide(prs, title, subtitle, *, deck_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, BG_HERO)
    add_text(slide, Inches(0.8), Inches(2.6), Inches(11), Inches(1.0),
             title, size=48, bold=True, color=INK_INVERT)
    add_text(slide, Inches(0.8), Inches(3.6), Inches(11), Inches(0.6),
             subtitle, size=22, color=INK_MUTED_INVERT)
    add_text(slide, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
             "Wayve  ·  dev.maheshg.me  ·  API version 2026.05",
             size=12, color=INK_MUTED_INVERT)
    return slide


def make_section_slide(prs, section_no, title, blurb, *, deck_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_NAV)
    add_text(slide, Inches(0.8), Inches(2.4), Inches(11), Inches(0.6),
             f"§ {section_no}", size=24, color=ACCENT_LIGHT)
    add_text(slide, Inches(0.8), Inches(3.0), Inches(11), Inches(1.4),
             title, size=42, bold=True, color=INK_INVERT)
    add_text(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(1.2),
             blurb, size=18, color=INK_MUTED_INVERT)
    return slide


def make_content_slide(prs, title, kicker=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_PANEL)
    if kicker:
        add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.4),
                 kicker, size=12, bold=True, color=ACCENT)
    add_text(slide, Inches(0.8), Inches(0.7), Inches(12), Inches(0.9),
             title, size=30, bold=True, color=INK_PRIMARY)
    # Underline accent
    add_rect(slide, Inches(0.8), Inches(1.55), Inches(0.8), Inches(0.06), ACCENT)
    return slide


def make_closer_slide(prs, title, blurb, *, deck_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_HERO)
    add_text(slide, Inches(0.8), Inches(2.6), Inches(11.5), Inches(1.4),
             title, size=42, bold=True, color=INK_INVERT)
    add_text(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(1.0),
             blurb, size=20, color=INK_MUTED_INVERT)
    add_text(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
             "Cross-reference: wayve.md  ·  technical_wayve.md  ·  business_wayve.md",
             size=12, color=INK_MUTED_INVERT)
    return slide


# ── Technical deck ────────────────────────────────────────────────────


def build_technical_deck(out_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    deck_name = "Technical Deep Dive"
    slides_built = []

    # 1. Title
    slides_built.append(make_title_slide(
        prs,
        "Wayve — Technical Deep Dive",
        "Architecture · auth · workers · webhooks · scaling characteristics",
        deck_name=deck_name,
    ))

    # 2. Agenda
    s = make_content_slide(prs, "Agenda", "01")
    add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
                [
                    "Stack & topology — Rust + React + Postgres + Redis + WS",
                    "Request lifecycle — JWT, X-API-KEY, X-EMBED-TOKEN, SCIM bearer",
                    "RBAC — 3 scopes × 9 roles × 22 permissions",
                    "Rate limit + quota — MIN(key, plan) + per-user monthly counter",
                    "Workers — sync, body backfill, dispatcher",
                    "Webhooks — signed deliveries with retry",
                    "Encryption — at-rest AES-GCM + E2E chat envelope",
                    "Threat model + scaling bottlenecks",
                ],
                size=18)
    slides_built.append(s)

    # 3. Section: architecture
    slides_built.append(make_section_slide(prs, "1", "Architecture overview",
                                            "How a request becomes a row.",
                                            deck_name=deck_name))

    # 4. Topology
    s = make_content_slide(prs, "Single-instance v1 topology", "02")
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.5),
             "EC2 i-07af9db286562f5ac  ·  EIP 32.199.117.86  ·  dev.maheshg.me",
             size=14, color=INK_MUTED)
    layers = [
        ("nginx 1.27", "TLS terminate · CSP · HSTS · routes /api,/scim/v2,/ws/* → backend; else SPA"),
        ("rwayve-backend (Rust)", "Actix Web 4 · 11 feature modules · 4 runtime roles via RWAYVE_ROLE"),
        ("Postgres 15 · Redis 7 · MailHog · Jaeger", "Postgres = source of truth; Redis = counters + cache; MailHog dev only; Jaeger OTLP"),
        ("External", "Gmail · MS Graph · Stripe · Gemini · Cloudflare TURN · SMTP"),
    ]
    for i, (heading, body) in enumerate(layers):
        y = Inches(2.6 + i * 1.05)
        add_rect(s, Inches(0.8), y, Inches(11.5), Inches(0.95), BG_PANEL, line=DIVIDER)
        add_text(s, Inches(1.0), y + Inches(0.08), Inches(11), Inches(0.4),
                 heading, size=16, bold=True, color=BG_NAV)
        add_text(s, Inches(1.0), y + Inches(0.45), Inches(11), Inches(0.5),
                 body, size=12, color=INK_MUTED)
    slides_built.append(s)

    # 5. Runtime roles
    s = make_content_slide(prs, "Four runtime roles, one binary", "03")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.5),
              ["RWAYVE_ROLE", "Spawned components", "DB pool"],
              [
                  ["api (default)", "HTTP server + webhook dispatcher", "10"],
                  ["email_sync_worker", "email::sync::sync_all loop", "5"],
                  ["email_body_worker", "body backfill (40 concurrent / account)", "5"],
                  ["all", "API + every worker (dev only)", "10"],
              ])
    add_text(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2),
             "Why split workers?  Sync is bursty (30 s tick → up to 50 messages × N accounts).\n"
             "Body backfill is concurrency-bound (40 parallel fetches per account).\n"
             "Sharing pool slots with the HTTP server would starve either one under load.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 6. Auth chain
    slides_built.append(make_section_slide(prs, "2", "Request lifecycle",
                                            "Middleware chain: API key → embed → CORS → handler.",
                                            deck_name=deck_name))

    # 7. Auth mechanisms
    s = make_content_slide(prs, "Four authentication mechanisms", "04")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5),
              ["Mechanism", "Header", "TTL", "Where it lands"],
              [
                  ["JWT session", "Authorization: Bearer … OR cookie rwayve_auth", "24 h", "Every endpoint"],
                  ["X-API-KEY", "X-API-KEY: wv_sk_…", "until revoked / expires", "/api/* (excluding management)"],
                  ["X-EMBED-TOKEN", "X-EMBED-TOKEN: <jwt>", "5 min", "/api/* GET/HEAD only, origin-pinned"],
                  ["SCIM bearer", "Authorization: Bearer wv_scim_…", "until revoked", "/scim/v2/* only"],
                  ["OIDC SSO", "(redirect-based; consumes session cookie after)", "session", "Login flow only"],
              ])
    slides_built.append(s)

    # 8. get_user_id_from_request
    s = make_content_slide(prs, "Resolver order matters", "05")
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.6),
             "backend/src/security/jwt.rs::get_user_id_from_request",
             size=14, color=INK_MUTED, font="Consolas")
    add_rect(s, Inches(0.8), Inches(2.6), Inches(11.5), Inches(3.6), BG_HERO)
    code = (
        "pub fn get_user_id_from_request(req: &HttpRequest) -> Option<i32> {\n"
        "    // 1. API-key principal (stamped by ApiKeyMiddleware)\n"
        "    if let Some(p) = req.extensions().get::<ApiKeyPrincipal>() {\n"
        "        return Some(p.user_id);\n"
        "    }\n"
        "    // 2. Embed-token principal (stamped by EmbedMiddleware)\n"
        "    if let Some(p) = req.extensions().get::<EmbedPrincipal>() {\n"
        "        return Some(p.user_id);\n"
        "    }\n"
        "    // 3. JWT (Authorization header OR cookie)\n"
        "    let token  = token_from_request(req)?;\n"
        "    let claims = decode_jwt(&token)?;\n"
        "    Some(claims.sub)\n"
        "}"
    )
    add_text(s, Inches(0.95), Inches(2.7), Inches(11.2), Inches(3.4),
             code, size=12, color=INK_INVERT, font="Consolas")
    add_text(s, Inches(0.8), Inches(6.3), Inches(11.5), Inches(0.8),
             "Ordering matters — an API-key request never reaches the JWT path.\n"
             "WS endpoints fall back to ?token= query param; still verified, never trusted blindly.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 9. Section RBAC
    slides_built.append(make_section_slide(prs, "3", "RBAC",
                                            "9 roles × 22 permissions, computed from the DB on every request.",
                                            deck_name=deck_name))

    # 10. RBAC matrix snippet
    s = make_content_slide(prs, "Selected role / permission matrix", "06")
    add_table(s, Inches(0.4), Inches(2.0), Inches(12.5), Inches(4),
              ["Permission", "owner", "super_admin", "admin", "security", "billing", "developer", "support"],
              [
                  ["members:manage", "✓", "✓", "✓", "—", "—", "—", "—"],
                  ["roles:manage", "✓", "—", "—", "—", "—", "—", "—"],
                  ["billing:manage", "✓", "—", "—", "—", "✓", "—", "—"],
                  ["api_keys:manage", "✓", "✓", "—", "—", "—", "✓", "—"],
                  ["webhooks:manage", "✓", "✓", "—", "✓", "—", "✓", "—"],
                  ["audit:read", "✓", "✓", "—", "✓", "—", "—", "—"],
                  ["security:manage", "✓", "✓", "—", "✓", "—", "—", "—"],
                  ["tickets:manage", "✓", "✓", "—", "—", "—", "—", "✓"],
              ], font_size=11)
    add_text(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.6),
             "Source of truth: backend/src/security/rbac.rs · Frontend mirror: frontend/src/auth/permissions.ts",
             size=12, color=INK_MUTED, font="Consolas")
    slides_built.append(s)

    # 11. RBAC resolution
    s = make_content_slide(prs, "Per-request resolution, never JWT-trusted", "07")
    add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
                [
                    "Every privileged handler calls rbac::require_permission(req, pool, Permission::X)",
                    "resolve_role_context() joins users + organization_members + platform_members",
                    "Cost: < 1 ms per request (PK + two indexed LEFT JOINs)",
                    "Role change at T+0 → demoted user's next request at T+1 sees new perms",
                    "JWT only proves identity; capability is data, not capability",
                    "Future cache: 60 s Moka TTL keyed on user_id (mirror email::account pattern)",
                ],
                size=18)
    slides_built.append(s)

    # 12. Section: rate limit
    slides_built.append(make_section_slide(prs, "4", "Rate limit + quota",
                                            "MIN(key, plan) per minute + per-user monthly counter, Redis-backed.",
                                            deck_name=deck_name))

    # 13. Rate limit flow
    s = make_content_slide(prs, "9-step middleware flow per request", "08")
    add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
                [
                    "1. Header present? otherwise passthrough",
                    "2. resolve_api_key → SHA-256 hash lookup; revoked/expired = 401",
                    "3. required_scope(method, path) → option<&str>; unmapped = deny",
                    "4. scope_satisfied(key.scopes, required) → 403 if no",
                    "5. quotas::effective_for_user(pool, user_id) → cached plan tier",
                    "6. effective = MIN(key.rate_limit, plan.rate_limit); INCR apikey_rl",
                    "7. Per-user monthly INCR apikey_quota:{user_id}:{YYYY-MM}",
                    "8. Inject ApiKeyPrincipal into request extensions",
                    "9. After handler returns: audit::write — fire-and-forget",
                ],
                size=15)
    slides_built.append(s)

    # 14. Failure modes
    s = make_content_slide(prs, "Failure modes worth knowing", "09")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5),
              ["Scenario", "Behaviour"],
              [
                  ["Redis down at request time", "Both counters fail open; audit log records the relax"],
                  ["Postgres down at audit-write", "Audit entry lost (tokio::spawn, not blocked)"],
                  ["Postgres down at validate", "Request blocked with 500; no key validation possible"],
                  ["Plan downgrade mid-month", "Up to 60s cache TTL before new lower cap kicks in"],
                  ["Plan upgrade mid-month", "Same 60s delay; symmetric"],
                  ["Monthly quota exhausted", "429 with 'Monthly request quota exceeded' until cycle reset"],
              ])
    slides_built.append(s)

    # 15. Tier table
    s = make_content_slide(prs, "Pricing → rate limit ladder", "10")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5),
              ["Tier", "Plan code", "Price", "Rate", "Monthly requests"],
              [
                  ["Free", "basic_user", "$0", "60 / min", "50,000"],
                  ["Advance", "advance_user", "$7 / mo", "300 / min", "500,000"],
                  ["Organization", "organization", "$10 / seat / mo", "600 / min", "5,000,000"],
                  ["Enterprise", "enterprise", "Custom", "6,000 / min", "Unlimited"],
              ])
    add_text(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1),
             "Stricter cap wins: MIN(api_key.rate_limit_per_min, plan.rate_limit_per_min).\n"
             "Monthly counter aggregates across every key the user owns — splitting keys does not multiply budget.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 16. Section: workers
    slides_built.append(make_section_slide(prs, "5", "Background workers",
                                            "Sync · body backfill · billing reconcile · webhook dispatcher.",
                                            deck_name=deck_name))

    # 17. Worker table
    s = make_content_slide(prs, "Four workers, four roles", "11")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5),
              ["Worker", "Cadence", "Concurrency", "Idempotency"],
              [
                  ["sync_worker", "30 s tick", "20 accounts in flight", "ON CONFLICT DO UPDATE; xmax=0 = real INSERT"],
                  ["body_worker", "continuous (5 s idle)", "40 per account, 200/iter", "partial index on body_encrypted=''"],
                  ["billing reconcile", "1 h", "1", "idempotent UPDATE entitlements"],
                  ["webhook dispatcher", "5 s poll", "25-row claim batch", "FOR UPDATE SKIP LOCKED"],
              ])
    add_text(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
             "On error: exp backoff up to 5 min (sync), claim retry on next tick (dispatcher).",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 18. Section: webhooks
    slides_built.append(make_section_slide(prs, "6", "Webhook delivery",
                                            "Signed envelopes, 3-attempt retry, auto-disable after 20 failures.",
                                            deck_name=deck_name))

    # 19. Envelope
    s = make_content_slide(prs, "Event envelope + signing", "12")
    add_rect(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.4), BG_HERO)
    envelope_code = (
        "{\n"
        '  "id":          "evt_<uuid>",\n'
        '  "type":        "task.created",\n'
        '  "api_version": "2026.05",\n'
        '  "created_at":  "2026-05-24T13:45:00Z",\n'
        '  "owner":       { "type": "user", "user_id": 42, "organization_id": null },\n'
        '  "data":        { /* event-specific payload */ }\n'
        "}\n"
        "\n"
        "HTTP/1.1 POST https://customer.example/wayve\n"
        "Wayve-Signature: t=1779629510,v1=<HMAC-SHA256(secret, t.body)>"
    )
    add_text(s, Inches(0.95), Inches(2.1), Inches(11.2), Inches(3.2),
             envelope_code, size=12, color=INK_INVERT, font="Consolas")
    add_text(s, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.4),
             "Replay defence: signature pinned to timestamp; receivers reject anything > 5 min old.\n"
             "Idempotency: receiver dedupes on evt_id. Wayve makes no exactly-once claim.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 20. Event catalog
    s = make_content_slide(prs, "11-event catalog (frozen v1)", "13")
    add_table(s, Inches(0.4), Inches(2.0), Inches(12.6), Inches(4.6),
              ["Event", "Fires from", "Payload notes"],
              [
                  ["task.created/.updated/.deleted", "tasks/handler.rs", "Full task row / {id}"],
                  ["meeting.created/.updated/.deleted", "scheduler/handler.rs", "Title + date + zoom URL + participants"],
                  ["email.received", "email/sync.rs (xmax=0 only)", "id, account_id, sender, subject. NOT body."],
                  ["email.sent", "email/send.rs (2xx only)", "from, to, subject. NOT body."],
                  ["chat.message.sent", "chat/websocket.rs", "Metadata only — content is E2E encrypted"],
                  ["chat.channel.created", "chat/channel_create.rs", "post-commit only"],
                  ["wayve.ping", "POST /webhooks/{id}/test", "Test verification"],
              ], font_size=11)
    slides_built.append(s)

    # 21. Section: encryption
    slides_built.append(make_section_slide(prs, "7", "Encryption",
                                            "AES-GCM at rest + E2E envelope for chat.",
                                            deck_name=deck_name))

    # 22. Key derivation
    s = make_content_slide(prs, "Layered model", "14")
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.5),
             "What Wayve actually owns is layer 1.",
             size=14, color=INK_MUTED)
    add_table(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(2.8),
              ["Layer", "Mechanism", "What it protects"],
              [
                  ["1. App-level", "AES-256-GCM via HKDF-SHA512(AES_KEY, salt)", "DB dump leak"],
                  ["2. Postgres", "text columns hold ciphertext", "—"],
                  ["3. EBS", "AWS KMS-managed encryption", "Disk seizure"],
              ])
    add_text(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
             "Nonce: 12 bytes random per encryption.\n"
             "Legacy fallback: tries raw AES_KEY if HKDF-derived fails — enables key rotation.\n"
             "AES_HKDF_SALT must be stable once set; rotating it loses all rows.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 23. E2E chat
    s = make_content_slide(prs, "E2E chat envelope", "15")
    add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
                [
                    "Client generates fresh 32-byte AES key K + 12-byte nonce per message",
                    "AES-GCM-encrypts plaintext → (iv_payload, ct_payload)",
                    "RSA-OAEP-wraps K with each recipient's public key",
                    "Envelope: WAYVE_CHAT_E2E_V1\\n{wrapped_keys, iv, ct}",
                    "Server refuses unwrapped messages at WS receive",
                    "Server stores envelope inside an at-rest AES layer (defense in depth)",
                    "Recipient browser unwraps using IndexedDB-stored private key",
                    "v1 gaps: no forward secrecy, no double ratchet, RSA-2048 (not Curve25519)",
                ],
                size=16)
    slides_built.append(s)

    # 24. Section: threat model
    slides_built.append(make_section_slide(prs, "8", "Threat model + scaling",
                                            "STRIDE notes and where the system breaks first.",
                                            deck_name=deck_name))

    # 25. STRIDE summary
    s = make_content_slide(prs, "STRIDE highlights", "16")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5),
              ["Category", "Mitigated", "Honest gap"],
              [
                  ["Spoofing", "JWT HS256 + API-key hash + webhook HMAC timestamp", "No JWT jti revocation list (24h leak window)"],
                  ["Tampering", "HTTPS + HMAC + AEAD tags", "—"],
                  ["Repudiation", "api_key_audit_log every request", "No login/role-change audit events yet"],
                  ["Info disclosure", "Encrypted at rest; tracing skips PII", "Side-channel via response timing not measured"],
                  ["DoS", "Rate limit + monthly quota", "Single-EC2 SPOF"],
                  ["Elevation", "Per-request RBAC resolve", "—"],
              ])
    slides_built.append(s)

    # 26. Scaling bottlenecks
    s = make_content_slide(prs, "Bottlenecks in order they bite", "17")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5),
              ["Bottleneck", "Trigger", "Fix"],
              [
                  ["Single Postgres", "Read q/s > 200", "Add read replica (~$50/mo)"],
                  ["Single EC2", "Saturated CPU/mem", "ECS Fargate fleet behind ALB"],
                  ["Webhook fan-out", "100 subscribers × 1 event = 20s drain", "Bump dispatcher batch; partition"],
                  ["Bcrypt on login", "~100ms / login CPU", "Per-IP login rate limit; argon2id later"],
                  ["Email sync ceiling", "> 20 concurrent accounts", "Shard accounts across worker containers"],
              ])
    slides_built.append(s)

    # 27. Performance numbers
    s = make_content_slide(prs, "Measured at v1 scale", "18")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5),
              ["Endpoint / op", "p95 latency"],
              [
                  ["GET /api/me", "~ 50 ms (1 PG + 1 Redis)"],
                  ["GET /api/emails (75 rows decrypted)", "~ 120 ms"],
                  ["GET /api/openapi.json (cached + ETag)", "~ 5 ms"],
                  ["WS handshake (chat or call)", "~ 80 ms (incl. RBAC)"],
                  ["Webhook delivery end-to-end", "1-2 s (5 s claim poll + HTTP RTT)"],
              ])
    add_text(s, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1),
             "Postgres + Redis on same EC2; no replica; mimalloc allocator.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 28. Tech debt
    s = make_content_slide(prs, "Known tech debt, ranked", "19")
    add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
                [
                    "Self-healing migrations → switch to sqlx migrate around 50 ALTERs",
                    "Single-EC2 SPOF — container fleet + ALB",
                    "Drive on EBS — move to S3 with signed-URL gateway",
                    "audit_events table missing for login + role change",
                    "No OAuth 2.0 authorization server (Wayve is RP only)",
                    "RSA-2048 chat envelope — Signal-grade audit would balk",
                    "No connection pooler (PgBouncer)",
                    "No metrics → no real SLOs (add Prometheus)",
                ],
                size=16)
    slides_built.append(s)

    # 29. Closer
    slides_built.append(make_closer_slide(
        prs,
        "Engineering is the contract",
        "We don't hide the gaps. OpenAPI + audit + known-limits docs are deliberate brand.",
        deck_name=deck_name,
    ))

    total = len(slides_built)
    for idx, sl in enumerate(slides_built, start=1):
        if idx in (1, total):
            continue
        add_footer(sl, idx, total, deck_name)

    prs.save(out_path)
    print(f"  ✓ {out_path}  ({total} slides)")


# ── Business deck ─────────────────────────────────────────────────────


def build_business_deck(out_path: Path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    deck_name = "Business Deep Dive"
    slides_built = []

    # 1. Title
    slides_built.append(make_title_slide(
        prs,
        "Wayve — Business Deep Dive",
        "Market · pricing · unit economics · GTM · roadmap · risk",
        deck_name=deck_name,
    ))

    # 2. Pitch in 3 sentences
    s = make_content_slide(prs, "Wayve in three sentences", "01")
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
             "Mid-market companies pay $50-100 per seat per month stitching together\n"
             "Gmail + Slack + Notion + Drive + Zoom + Calendly.\n\n"
             "Wayve replaces that bundle for $10 per seat per month — with end-to-end\n"
             "chat encryption, SCIM, audit export, and a self-hostable binary —\n"
             "the enterprise primitives the bundled megacorps either don't ship\n"
             "or hide behind a $50/seat Enterprise tier.",
             size=22, color=INK_PRIMARY)
    slides_built.append(s)

    # 3. The problem
    s = make_content_slide(prs, "The problem", "02")
    add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
                [
                    "Mid-market typically runs 12-15 SaaS tools per knowledge worker",
                    "Per-seat cost across the stack: $50-100 / month / employee",
                    "Five SCIM configs · five access reviews · five DPAs · five invoices",
                    "SOC 2 + GDPR procurement makes vendor sprawl an enumerable cost",
                    "No tool talks to the others — integrations are an ongoing engineering tax",
                ],
                size=20)
    slides_built.append(s)

    # 4. The solution
    s = make_content_slide(prs, "The solution", "03")
    add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
                [
                    "One workspace: email, chat, calls, scheduler, drive, notes, tasks, AI",
                    "$10 / seat / month — one bill, one SSO, one audit log, one API",
                    "End-to-end encryption on chat (envelope-layer, server cannot read)",
                    "Programmatic API on day one: OpenAPI spec + dev portal + signed webhooks",
                    "Self-hostable: same binary serves SaaS and on-prem (EU sovereignty win)",
                ],
                size=20)
    slides_built.append(s)

    # 5. Why now
    s = make_content_slide(prs, "Why now", "04")
    add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
                [
                    "AI ratchet: every productivity tool just added LLM features. Winners will be the suites with workspace context, not point tools.",
                    "Privacy regulation tailwind: Schrems II, GDPR rulings, EU AI Act all push for data-sovereign alternatives to US hyperscalers.",
                    "Tool fatigue: customer-success leaders report 12-15 SaaS tools per employee. Procurement reviews are hammering this.",
                    "Mid-market gap: Microsoft 365 + Google Workspace either too expensive (Enterprise SKUs) or missing the encryption story.",
                ],
                size=18)
    slides_built.append(s)

    # 6. Section: market
    slides_built.append(make_section_slide(prs, "1", "Market", "TAM, SAM, SOM — sized by knowledge-worker seat-month.",
                                            deck_name=deck_name))

    # 7. TAM/SAM/SOM
    s = make_content_slide(prs, "$30 B TAM, $25 B SAM, $360k year-1 SOM", "05")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5),
              ["Layer", "Math", "Annual"],
              [
                  ["TAM", "250 M knowledge-worker seats × $10 ARPU/mo × 12", "$30 B"],
                  ["SAM", "210 M seats in addressable mid-market", "$25 B"],
                  ["SOM (Year 1, conservative)", "100 customers × 30 seats × $10 × 12", "$360 k"],
                  ["SOM (Year 3, realistic)", "1,500 customers × 40 seats × $10 × 12", "$7.2 M"],
              ])
    add_text(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8),
             "TAM = ceiling, not forecast. Microsoft Teams alone is estimated at $5-8B within Office 365.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 8. Section: competition
    slides_built.append(make_section_slide(prs, "2", "Competition", "Bundled suites and point tools — we sit in between.",
                                            deck_name=deck_name))

    # 9. Competition table
    s = make_content_slide(prs, "Wayve vs the bundled suites", "06")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4),
              ["Product", "Bundle?", "Per seat / mo", "E2E chat?", "Self-host?"],
              [
                  ["Wayve Organization", "7 apps", "$10", "✓", "✓"],
                  ["Google Workspace Biz+", "5 apps", "$18", "—", "—"],
                  ["Microsoft 365 Biz Prem", "5 apps + Office", "$22", "—", "—"],
                  ["Slack Business+", "chat only", "$15", "—", "—"],
                  ["Notion Business", "docs only", "$15", "—", "—"],
                  ["Zoho One", "40+ apps", "$37", "—", "✓ (some)"],
              ])
    add_text(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
             "Stitched alternative for an equivalent stack: $50-100 across multiple invoices.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 10. The bet
    s = make_content_slide(prs, "The 'good enough at each, better at the bundle' bet", "07")
    add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
                [
                    "Wayve will NOT beat Slack at search or Notion at databases for 2-3 years",
                    "80% of mid-market knowledge work doesn't need best-in-class for any single product",
                    "It needs coherent enough across all of them: one bill, one SSO, one API",
                    "Precedent: HubSpot vs Salesforce+Pardot+ZoomInfo. Atlassian beat the point tools.",
                    "Our structural advantage: API + encryption + self-host on day one, not retrofitted",
                ],
                size=18)
    slides_built.append(s)

    # 11. Section: ICP
    slides_built.append(make_section_slide(prs, "3", "Ideal customer", "Mid-market privacy-conscious tech companies, 20-200 seats.",
                                            deck_name=deck_name))

    # 12. ICP table
    s = make_content_slide(prs, "Three customer segments", "08")
    add_table(s, Inches(0.4), Inches(2.0), Inches(12.6), Inches(4.5),
              ["Segment", "Example", "Champion", "ACV", "Cycle"],
              [
                  ["Tier-1: privacy-conscious tech", "Privacy-first SaaS, infosec, legal-tech, EU AI", "VPE / CTO / IT lead", "$6,000", "2-6 weeks"],
                  ["Tier-2: sensitive projects at larger Cos", "Defense, M&A, finance MNPI", "Project lead with budget", "$5k - $50k", "2-4 weeks"],
                  ["Tier-3: personal power users", "Individual developers", "Self-onboarding", "$0-84 / yr", "—"],
              ], font_size=11)
    add_text(s, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.5),
             "Anti-ICP: Fortune 500 (procurement = 18mo), consumer, single-product replacements.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 13. Section: pricing
    slides_built.append(make_section_slide(prs, "4", "Pricing strategy", "Four tiers, anchored on the $10 / seat / month Organization plan.",
                                            deck_name=deck_name))

    # 14. Pricing tiers
    s = make_content_slide(prs, "Four tiers, ground truth from the codebase", "09")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4),
              ["Tier", "Price", "Rate / min", "Monthly API", "Audience"],
              [
                  ["Free", "$0", "60", "50,000", "Personal eval / OSS contributors"],
                  ["Advance", "$7 / mo", "300", "500,000", "Hobby builders, freelancers"],
                  ["Organization", "$10 / seat / mo", "600", "5,000,000", "10-500 seat companies (revenue tier)"],
                  ["Enterprise", "Custom", "6,000", "Unlimited", "500+ seats, SLA + on-prem"],
              ])
    add_text(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6),
             "Free tier exists to seed developer integrations — not as a consumer product.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 15. Section: unit economics
    slides_built.append(make_section_slide(prs, "5", "Unit economics", "$6k ACV · 88% gross margin · 36-month customer lifetime.",
                                            deck_name=deck_name))

    # 16. Unit econ
    s = make_content_slide(prs, "Per-customer math, Organization tier", "10")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5),
              ["Metric", "Value"],
              [
                  ["Average seats", "50"],
                  ["ACV", "$6,000 / year"],
                  ["Customer lifetime", "36 months"],
                  ["LTV", "$18,000"],
                  ["Variable COGS / seat / mo", "$1.15"],
                  ["Gross margin", "88.5%"],
                  ["LTV (gross)", "$14,400"],
              ])
    add_text(s, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.2),
             "CAC target year 1: ~$400 blended (60% PLG + 40% founder-led outbound).\n"
             "LTV:CAC = 36:1 today (optimistic); ~14:1 once outbound matures (still excellent).\n"
             "CAC payback: 2.4 months.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 17. Section: GTM
    slides_built.append(make_section_slide(prs, "6", "Go-to-market", "PLG + founder-led outbound + dev-portal pull.",
                                            deck_name=deck_name))

    # 18. Funnel
    s = make_content_slide(prs, "The funnel", "11")
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
             "1. SEO + content       → comparison pages, eng blog\n"
             "2. Developer portal    → /developers + tutorials drives evals\n"
             "3. Free signup         → individual API + UI access\n"
             "4. Bring to team       → champion invites colleagues\n"
             "5. Org-tier upgrade    → admin upgrades workspace\n"
             "6. Expansion           → more seats + Enterprise SLA",
             size=22, color=INK_PRIMARY, font="Consolas")
    slides_built.append(s)

    # 19. Channel weighting
    s = make_content_slide(prs, "Channel mix by year", "12")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5),
              ["Channel", "Year 1", "Year 2", "Year 3"],
              [
                  ["Organic search + content", "40%", "25%", "15%"],
                  ["Developer-portal funnel", "25%", "30%", "30%"],
                  ["Founder-led outbound", "30%", "15%", "5%"],
                  ["SDR-led outbound", "0%", "15%", "25%"],
                  ["Partner / referral", "5%", "15%", "25%"],
              ])
    add_text(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.8),
             "Founder-led outbound caps at ~50 logos / year. SDR-led requires hiring sales.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 20. Section: revenue projections
    slides_built.append(make_section_slide(prs, "7", "Revenue", "Three projections — none are forecasts, all are stress tests.",
                                            deck_name=deck_name))

    # 21. Realistic projection
    s = make_content_slide(prs, "Realistic projection", "13")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4),
              ["Period", "Customers", "Avg seats", "MRR", "ARR"],
              [
                  ["Month 6", "40", "25", "$10k", "$120k"],
                  ["Month 12", "150", "30", "$45k", "$540k"],
                  ["Month 18", "380", "32", "$121.6k", "$1.46M"],
                  ["Month 24", "720", "35", "$252k", "$3.02M"],
                  ["Month 36", "1,500", "40", "$600k", "$7.2M"],
              ])
    add_text(s, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.6),
             "Dominant lever: new-logo acquisition rate. Keep > 12% MoM in years 1-2.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 22. Revenue mix year 3
    s = make_content_slide(prs, "Revenue mix at year 3 (realistic)", "14")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3),
              ["Tier", "% of customers", "% of ARR"],
              [
                  ["Free", "60% (developer base)", "0%"],
                  ["Advance ($7/mo)", "25%", "3%"],
                  ["Organization ($10/seat)", "14%", "75%"],
                  ["Enterprise (custom)", "1%", "22%"],
              ])
    add_text(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.5),
             "Enterprise concentration (22%) requires landing 5-10 enterprise logos by Year 3,\n"
             "each contributing $100k-$500k ARR. Without enterprise, Year 3 ARR shape shifts:\n"
             "Org alone tops out around ~$5.5M ARR at the same logo count.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 23. Section: roadmap
    slides_built.append(make_section_slide(prs, "8", "Roadmap", "12 months: bundle polish + AI + SOC 2. 36 months: deep enterprise.",
                                            deck_name=deck_name))

    # 24. 12-month roadmap
    s = make_content_slide(prs, "12-month roadmap (priorities ordered)", "15")
    add_bullets(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
                [
                    "Universal search — single biggest 'wow' the bundle can deliver",
                    "AI workspace context — Gemini bridge with real cross-product visibility",
                    "SOC 2 Type II — procurement blocker for 60%+ of mid-market deals",
                    "Mobile apps — iOS + Android (required for chat reliability)",
                    "Calendar polish — recurring meetings, time zones, sharing UX",
                    "Audit retention tiers — 30/90/365d as pricing axis",
                    "Wayve as OAuth 2.0 authorization server (when first partner asks)",
                ],
                size=18)
    slides_built.append(s)

    # 25. Section: risk
    slides_built.append(make_section_slide(prs, "9", "Risk register", "Bus factor first. Then compliance. Then technical.",
                                            deck_name=deck_name))

    # 26. Top risks
    s = make_content_slide(prs, "Top 6 risks, ranked", "16")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5),
              ["Risk", "Likelihood", "Impact", "Mitigation"],
              [
                  ["Bus factor (single founder)", "High", "Existential", "Hire 2nd founding eng by month 9"],
                  ["MSFT/Google bundles a similar SKU", "Medium", "High", "Lean harder into encryption + self-host"],
                  ["SOC 2 Type II doesn't land Yr 1", "Medium", "Blocks 60% deals", "Engage auditor immediately"],
                  ["Email sync unreliable under load", "Medium", "Churn", "Shard accounts across workers"],
                  ["AWS bill surprise (10M API/day)", "Low", "Operational", "Quota system in place"],
                  ["Founder burnout", "High", "Existential", "Hire by month 9; no marathons past Q2"],
              ])
    slides_built.append(s)

    # 27. Section: compliance
    slides_built.append(make_section_slide(prs, "10", "Compliance", "GDPR ready today. SOC 2 Type II Year 1. HIPAA Year 2.",
                                            deck_name=deck_name))

    # 28. Compliance posture
    s = make_content_slide(prs, "Compliance roadmap", "17")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4),
              ["Cert", "Status", "Cost", "Timeline"],
              [
                  ["GDPR", "Shipped today", "DPA template ~ 1 day", "Now"],
                  ["SOC 2 Type I", "Plan", "$15-30k", "Months 1-6"],
                  ["SOC 2 Type II", "Plan", "$30-60k total", "Months 1-13"],
                  ["HIPAA BAA", "Plan", "Add to Type II controls", "Month 18"],
                  ["ISO 27001", "Stretch", "$50-100k", "Year 2-3"],
                  ["FedRAMP", "No", "$2M+", "Not in plan"],
              ])
    slides_built.append(s)

    # 29. Section: team + funding
    slides_built.append(make_section_slide(prs, "11", "Team + funding", "Path A→B→C: bootstrap, optional pre-seed, seed at $50k MRR.",
                                            deck_name=deck_name))

    # 30. Hires
    s = make_content_slide(prs, "12-month hire plan", "18")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4),
              ["Role", "Type", "When", "Why"],
              [
                  ["Co-founder / 2nd engineer", "FT, equity", "Month 3-6", "Bus factor"],
                  ["Frontend engineer", "FT", "Month 6-9", "Mobile + UX polish"],
                  ["Customer support", "PT → FT", "Month 6", "Volume past founder's threshold"],
                  ["Designer", "Contractor → FT", "Month 9", "Brand for enterprise sales"],
                  ["Sales lead / SDR", "FT", "Month 12+", "Outbound motion"],
              ])
    slides_built.append(s)

    # 31. Funding path
    s = make_content_slide(prs, "Funding strategy", "19")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3),
              ["Path", "Trigger", "Size", "Dilution"],
              [
                  ["A. Bootstrap", "Day 0", "—", "0%"],
                  ["B. Pre-seed (SAFE)", "Month 3-6", "$150-500k", "5-10%"],
                  ["C. Seed round", "$50k MRR, 10% MoM growth", "$1.5-3M", "15-22%"],
              ])
    add_text(s, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.5),
             "Recommended path: A → B → C.\n"
             "Anti-pattern: raise too much too early — $5M seed forces growth-stage milestones.\n"
             "Better: smaller round, hit milestones, raise at higher valuation later.",
             size=14, color=INK_MUTED)
    slides_built.append(s)

    # 32. KPIs
    s = make_content_slide(prs, "North-star metrics", "20")
    add_table(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4),
              ["KPI", "Year-1 target"],
              [
                  ["MRR growth (MoM)", "10-15%"],
                  ["New logos / month", "20+ by month 12"],
                  ["Free → paid conversion", "5-8% of activated"],
                  ["Net revenue retention", "> 100%"],
                  ["Gross margin", "> 85%"],
                  ["p95 API latency", "< 200 ms"],
                  ["Uptime", "> 99.5%"],
                  ["Webhook delivery success", "> 99%"],
              ])
    slides_built.append(s)

    # 33. The ask
    s = make_content_slide(prs, "The ask", "21")
    add_text(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(5),
             "Raising $1.5 – $3M at the seed stage.\n\n"
             "Use of funds:\n"
             "   ─ 60% engineering hires (co-founder + 2 engineers)\n"
             "   ─ 20% compliance (SOC 2 Type II, HIPAA prep)\n"
             "   ─ 10% sales (1 SDR + tooling)\n"
             "   ─ 10% reserve\n\n"
             "Triggers: $50k MRR, ~10% MoM growth, 30+ paying logos, SOC 2 Type I in flight.",
             size=20, color=INK_PRIMARY)
    slides_built.append(s)

    # 34. Closer
    slides_built.append(make_closer_slide(
        prs,
        "Bundle. Encrypt. Honest pricing.",
        "Wayve · dev.maheshg.me · founder@dev.maheshg.me",
        deck_name=deck_name,
    ))

    total = len(slides_built)
    for idx, sl in enumerate(slides_built, start=1):
        if idx in (1, total):
            continue
        add_footer(sl, idx, total, deck_name)

    prs.save(out_path)
    print(f"  ✓ {out_path}  ({total} slides)")


# ── Entry ─────────────────────────────────────────────────────────────


def main():
    repo_root = Path(__file__).resolve().parent.parent
    print("Generating PowerPoint decks…")
    build_technical_deck(repo_root / "technical_wayve.pptx")
    build_business_deck(repo_root / "business_wayve.pptx")
    print("Done.")


if __name__ == "__main__":
    main()
